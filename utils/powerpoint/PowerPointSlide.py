
import numpy as np
import pandas as pd
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from apps.backends.LumoReporting.utils.powerpoint.PowerPointBase import BasePresentation


class BaseSlide:
    def __init__(self, presentation: BasePresentation = None, slide_name: str = None, data: dict = None):
        self.presentation = presentation
        self.slide_name = slide_name
        self.slide = self.presentation.add_slide(self.slide_name)
        self.data = data
        self.set_placeholders()

    def set_placeholders(self):
        self.placeholder_dict = {}
        for shape_name in ['Table Placeholder', 'Title', 'Text Placeholder', 'Chart Placeholder']:
            idx = 1
            for shape in self.slide.placeholders:
                if shape_name in shape.name:
                    shape.name = shape_name + ' ' + str(idx)
                    idx += 1
                    # print(shape.name, shape.placeholder_format.idx, shape.text_frame, shape.text)
                    self.placeholder_dict[shape.name] = shape.placeholder_format.idx
    def set_master_style(self, shape_text: str = None, replace_text: str = None, font_size: int = 6, right_alignment: bool = False ):
        for shape in self.presentation.presentation.slide_master.shapes:
            if not shape.has_text_frame:
                continue
            if shape.text == shape_text:
                shape.text = replace_text
                shape.text_frame.paragraphs[0].font.size = Pt(font_size)
                if right_alignment:
                    shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

    def get_placeholder(self, shape_name: str = None, shape_number: int = None):
        for key, item in self.placeholder_dict.items():
            if key == shape_name + ' ' + str(shape_number):
                placeholder = self.slide.placeholders[item]
                return placeholder
        raise ValueError(f'The number {shape_number} for the type {shape_name} does not exist. Edit template')

    def remove_placeholder(self, shape_name: str = None, shape_number: int = None) -> None:
        placeholder = self.get_placeholder(shape_name=shape_name, shape_number=shape_number)
        sp = placeholder.element
        sp.getparent().remove(sp)
    def set_superscript(self,
                        text_frame,
                        text: str = '',
                        supertext: str = '',
                        text_size=Pt(8),
                        super_size=Pt(8),
                        baseline: str = '30000',
                        add_text: str = '',
                        bold: bool = False):
        text_frame = text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = text_size
        run.font.bold = bold
        super_text = p.add_run()
        super_text.text = supertext
        super_text.font.size = super_size
        super_text.font.bold = bold
        super_text.font._element.set('baseline', baseline)
        additional_text = p.add_run()
        additional_text.text = add_text
        additional_text.font.size = text_size
        additional_text.font.bold = bold
        # if supertext.isdigit():
        #     self.presentation.update_footers()

    def set_hyperlink(self,
                      text_frame,
                      text: str = '',
                      text_size=Pt(8),
                      hyper_text: str = '',
                      hyper_link = None,
                      hyper_size=Pt(8),
                      bold: bool = False
                      ):
        text_frame = text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = text_size
        run.font.bold = bold
        hrun = p.add_run()
        hrun.text = hyper_text
        hrun.font.size = hyper_size
        hrun.font.bold = bold
        hlink = hrun.hyperlink
        hlink.address = hyper_link

    def sub_element(self,
                    parent,
                    tagname,
                    **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def set_cell_border(self,
                        cell,
                        border_color='423D36',
                        border_width='12700'):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        # if Right a:lnR, Left a:lnL, Top a:lnT
        lnR = self.sub_element(tcPr, 'a:lnB', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = self.sub_element(lnR, 'a:solidFill')
        srgbClr = self.sub_element(solidFill, 'a:srgbClr', val=border_color)
        return cell

    def set_col_widths(self,
                       cols,
                       widths: list = None):

        for idx, obj in enumerate(cols):
            obj.width = widths[idx]

    def set_row_heights(self,
                        rows,
                        heights: list = None):

        for idx, obj in enumerate(rows):
            obj.height = heights[idx]

    def set_none_cell_margins(self,
                              cell):
        cell.margin_bottom = Inches(0.0)
        cell.margin_top = Inches(0.0)
        cell.margin_left = Inches(0.0)
        cell.margin_right = Inches(0.0)

    def set_text(self, text: str = '', text_number: int = None, font_size: int = None):
        text_placeholder = self.get_placeholder(shape_name='Text Placeholder', shape_number=text_number)
        text_placeholder.text = text
        if font_size is not None:
            text_placeholder.text_frame.paragraphs[0].font.size = Pt(font_size)

    def set_table_title(self,
                        table,
                        df,
                        text: str = '',
                        super_text: str = '',
                        add_text: str = '',
                        title_size=Pt(12)):

        self.set_superscript(text_frame=table.cell(0, 0).text_frame,
                             text=text,
                             text_size=title_size,
                             supertext=super_text,
                             super_size=title_size,
                             add_text=add_text,
                             bold=False)

        table.cell(0, 0).merge(table.cell(0, df.shape[1] - 1))
        table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(11)
        self.set_none_cell_margins(table.cell(0, 0))
        table.cell(0, 0).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    def set_row(self,
                table,
                row_num: int = 1,
                row_vals: list = None,
                super_text: list = None,
                add_text: list = None,
                bold: bool = True,
                vertical_position=MSO_ANCHOR.BOTTOM,
                right_align: bool = True):

        if not super_text:
            super_text = [''] * len(row_vals)
        if not add_text:
            add_text = [''] * len(row_vals)

        for col_index, col_val in enumerate(row_vals):
            self.set_none_cell_margins(table.cell(row_num, col_index))
            self.set_superscript(text_frame=table.cell(row_num, col_index).text_frame,
                                 text=col_val,
                                 supertext=super_text[col_index],
                                 add_text=add_text[col_index],
                                 bold=bold)
            self.set_cell_border(table.cell(row_num, col_index), border_width='8000')
            table.cell(row_num, col_index).vertical_anchor = vertical_position
            if right_align:
                if col_index == 0:
                    table.cell(row_num, col_index).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                else:
                    table.cell(row_num, col_index).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

    def set_col(self,
                table,
                col_num: int = 0,
                col_vals: list = None,
                super_text: list = None,
                add_text: list = None,
                start_row: int = 2,
                bold: bool = True):

        if not super_text:
            super_text = [''] * len(col_vals)
        if not add_text:
            add_text = [''] * len(col_vals)

        for row_index, row_val in enumerate(col_vals):
            self.set_superscript(text_frame=table.cell(row_index + start_row, col_num).text_frame,
                                 text=row_val,
                                 supertext=super_text[row_index],
                                 add_text=add_text[row_index],
                                 bold=bold)

    def set_table_values(self,
                         table,
                         df,
                         title: bool = True,
                         col_widths: list = None,
                         row_heights: list = None,
                         right_align: bool = True):

        rows, cols = df.shape
        m = df.values
        if title:
            add_row = 2
        else:
            add_row = 1

        for row in range(rows):
            for col in range(cols):
                val = m[row, col]
                text = str(val)
                if isinstance(val, float):
                    if np.isnan(val):
                        text = ""
                    else:
                        text = str(round(val, 2))
                table.cell(row + add_row, col).text = text
                table.cell(row + add_row, col).text_frame.paragraphs[0].font.size = Pt(8)
                self.set_none_cell_margins(table.cell(row + add_row, col))
                table.cell(row + add_row, col).vertical_anchor = MSO_ANCHOR.MIDDLE
                self.set_cell_border(table.cell(row + add_row, col), border_width='4000')
                if right_align:
                    if col == 0:
                        table.cell(row + add_row, col).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    else:
                        table.cell(row + add_row, col).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
        if col_widths is not None:
            self.set_col_widths(cols=table.columns, widths=col_widths)
        if row_heights is not None:
            self.set_row_heights(rows=table.rows, heights=row_heights)

    def set_table(self,
                  table_number: int = None,
                  df: pd.DataFrame = None,
                  title: bool = True):
        table_placeholder = self.get_placeholder(shape_name='Table Placeholder', shape_number=table_number)
        rows, cols = df.shape
        if title:
            add_rows = 2
        else:
            add_rows = 1

        ph_graphic_frame = table_placeholder.insert_table(rows=rows + add_rows, cols=cols)
        table = ph_graphic_frame.table
        style_id = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
        tbl = ph_graphic_frame._element.graphic.graphicData.tbl
        tbl[0][-1].text = style_id

        return table

    def insert_table(self,
                     table_number,
                     df,
                     has_title: bool = True,
                     title: str = '',
                     super_title: str = '',
                     add_title: str = '',
                     headers: list = None,
                     super_headers: list = None,
                     add_headers: list = None,
                     vertical_pos=MSO_ANCHOR.BOTTOM,
                     right_align: bool = True,
                     col_widths: list = None,
                     row_heights: list = None):

        if headers is None:
            headers = list(df.columns)
            super_headers = ['']*len(headers)
            add_headers = ['']*len(headers)

        table = self.set_table(table_number=table_number, df=df)

        self.set_table_title(table=table,
                        df=df,
                        text=title,
                        super_text=super_title,
                        add_text=add_title)

        if has_title:
            row = 1
        else:
            row = 0
        self.set_row(table=table,
                     row_num=row,
                     row_vals=headers,
                     super_text=super_headers,
                     add_text=add_headers,
                     bold=True,
                     vertical_position=vertical_pos,
                     right_align=right_align)

        self.set_table_values(table=table,
                         df=df,
                         col_widths=col_widths,
                         row_heights=row_heights,
                         right_align=right_align)

        return table

    def insert_chart(self,
                     data: pd.DataFrame = None,
                     category: str = '',
                     chart_series: dict = None,
                     chart_num: int = None,
                     title: str = None,
                     chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED):
        data = data.replace(np.nan, None)
        data_series = data.to_dict('series')
        category = data_series.get(category)

        chart = CategoryChartData()

        chart.categories = category
        for series_name, column_name in chart_series.items():
            chart.add_series(series_name, data_series.get(column_name))

        graph_placeholder = self.get_placeholder(shape_name='Chart Placeholder', shape_number=chart_num)
        graphic_frame = graph_placeholder.insert_chart(chart_type, chart)
        chart = graphic_frame.chart

        chart.has_legend = True
        chart.has_label = True
        if title is not None:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
        else:
            chart.has_title = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(8)


        return chart
    def AttributeSheet(self):
        return None


