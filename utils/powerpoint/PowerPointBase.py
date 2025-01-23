
from pptx import Presentation


class BasePresentation:
    def __init__(
        self, file_path: str = None
    ):
        self.presentation = Presentation(file_path)
        self.counter = '1'
    def add_slide(self, slide_name: str = None):
        for idx, t in enumerate(self.presentation.slide_layouts):
            if t.name == slide_name:
                return self.presentation.slides.add_slide(self.presentation.slide_layouts[idx])

    def update_footers(self, add: int = 1):
        counter = int(self.counter) + add
        self.counter = str(counter)

    def save(self, output):
        self.presentation.save(output)

